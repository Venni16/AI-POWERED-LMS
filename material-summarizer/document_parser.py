import os
import logging
from typing import Optional
import pdfplumber
from docx import Document
import PyPDF2
from pptx import Presentation

logger = logging.getLogger(__name__)

def parse_document(file_path: str, file_extension: str) -> str:
    """
    Parse different document formats and extract text
    """
    try:
        if file_extension == '.pdf':
            return parse_pdf(file_path)
        elif file_extension in ['.docx', '.doc']:
            return parse_docx(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            return parse_pptx(file_path)
        elif file_extension == '.txt':
            return parse_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")
    except Exception as e:
        logger.error(f"Error parsing document {file_path}: {e}")
        raise

def parse_pdf(file_path: str) -> str:
    """
    Extract text from PDF using multiple methods for better coverage
    """
    text = ""
    
    # Method 1: Use pdfplumber (better for text-based PDFs)
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        logger.warning(f"pdfplumber failed: {e}")
    
    # Method 2: Use PyPDF2 as fallback
    if not text.strip():
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            logger.warning(f"PyPDF2 failed: {e}")
    
    if not text.strip():
        raise ValueError("Could not extract text from PDF")
    
    return clean_text(text)

def parse_docx(file_path: str) -> str:
    """
    Extract text from DOCX/DOC files
    """
    try:
        doc = Document(file_path)
        text = ""
        
        # Extract paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + "\n"
        
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error parsing DOCX file: {e}")
        raise

def parse_pptx(file_path: str) -> str:
    """
    Extract text from PowerPoint files
    """
    try:
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error parsing PPTX file: {e}")
        raise

def parse_txt(file_path: str) -> str:
    """
    Extract text from plain text files
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return clean_text(text)
    except UnicodeDecodeError:
        # Try different encodings
        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    text = file.read()
                return clean_text(text)
            except UnicodeDecodeError:
                continue
        raise ValueError("Could not decode text file with any encoding")

def clean_text(text: str) -> str:
    """
    Clean and normalize extracted text
    """
    # Remove excessive whitespace
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    
    # Remove very short lines that are likely formatting artifacts
    meaningful_lines = [line for line in lines if len(line) > 2]
    
    # Join with proper spacing
    cleaned_text = '\n'.join(meaningful_lines)
    
    # Remove multiple consecutive newlines
    while '\n\n\n' in cleaned_text:
        cleaned_text = cleaned_text.replace('\n\n\n', '\n\n')
    
    return cleaned_text.strip()